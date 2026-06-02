from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

OUTPUT_PATH = r"C:\Users\vfournand\Documents\socioscope-website\frontend\certificates_group2.pptx"
LOGO_PATH   = r"C:\Users\vfournand\Documents\socioscope-website\frontend\public\logo.png"
BASE_URL    = "https://thesocioscope.github.io/frontend/projects/"

# Brand colours
GREEN      = RGBColor(0x56, 0xB3, 0x40)   # vibrant green (header bg & lines)
DARK_GREEN = RGBColor(0x1A, 0x5C, 0x1A)   # deep green (titles)
BODY_CLR   = RGBColor(0x33, 0x33, 0x33)   # dark-grey body text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

W, H = 11.69, 8.27   # A4 landscape (inches)

initiatives = [
    {"code": "KE-083", "name": "Ubuntu Kreative Village",             "slug": "ubuntu-kreative-village"},
    {"code": "ZA-015", "name": "Khulisa Food and Nutrition NPO",      "slug": "khulisa-food-and-nutrition-npo"},
    {"code": "MA-005", "name": "Zitouna Life",                        "slug": "zitouna-life"},
    {"code": "TZ-005", "name": "Drive Inn Group",                     "slug": "drive-inn-group"},
    {"code": "RW-017", "name": "Uruhimbi Kageyo Cooperative",         "slug": "uruhimbi-kageyo-cooperative"},
    {"code": "IN-005", "name": "No Food Waste",                       "slug": "no-food-waste"},
    {"code": "LK-007", "name": "Seashore Garden",                     "slug": "seashore-garden"},
    {"code": "HK-008", "name": "Yau Tam Mei Holiday Camp Limited",    "slug": "yau-tam-mei-holiday-camp-limited"},
    {"code": "PY-001", "name": "Mboja’o",                        "slug": "mboja-o"},
    {"code": "CA-003", "name": "Collectif Récolte",              "slug": "collectif-recolte"},
    {"code": "US-018", "name": "Mad Agriculture",                     "slug": "mad-agriculture"},
    {"code": "GB-005", "name": "Loop:Frome",                          "slug": "loop-frome"},
    {"code": "FR-102", "name": "La Ruche qui dit Oui",                "slug": "la-ruche-qui-dit-oui"},
    {"code": "ES-007", "name": "Pescartes",                           "slug": "pescartes"},
    {"code": "DK-075", "name": "Bowline",                             "slug": "bowline"},
    {"code": "CH-019", "name": "Eaternity",                           "slug": "eaternity"},
    {"code": "PL-005", "name": "Eko Spichlerz",                       "slug": "eko-spichlerz"},
    {"code": "GE-010", "name": "Parki ar minda",                      "slug": "parki-ar-minda"},
    {"code": "MX-002", "name": "De la Chinampa",                      "slug": "de-la-chinampa"},
    {"code": "CR-010", "name": "CoopeTarrazú RL",                "slug": "coopetarrazu-rl"},
    {"code": "EC-004", "name": "Red de Guardianes de Semillas",       "slug": "red-de-guardianes-de-semillas"},
    {"code": "PE-012", "name": "Cirkula",                             "slug": "cirkula"},
    {"code": "AR-079", "name": "Colonia Jaime",                       "slug": "colonia-jaime"},
    {"code": "CO-019", "name": "AGROSAVIA",                           "slug": "corporacion-colombiana-de-investigacion-agropecuaria-agrosavia"},
    {"code": "NZ-002", "name": "TE KAKANO",                           "slug": "te-kakano"},
]

# ── helpers ──────────────────────────────────────────────────────────────────

def fnt(run, size, bold=False, italic=False, color=None, face="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = face
    if color:
        run.font.color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, no_line=True):
    sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if no_line:
        sh.line.fill.background()
    return sh

def add_hline(slide, x, y, w, color=GREEN):
    return add_rect(slide, x, y, w, 0.015, color)

def add_tb(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = wrap
    return tb

def set_char_spacing(run, centipoints):
    """Add letter spacing via raw XML (hundredths of a point)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(centipoints))

# ── slide factory ─────────────────────────────────────────────────────────────

def make_slide(prs, init):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    name = init["name"]
    url  = BASE_URL + init["slug"]

    # ── green header block (top-left) ────────────────────────────────────────
    GW, GH = 3.85, 2.35
    add_rect(slide, 0, 0, GW, GH, GREEN)

    # White triangle that cuts the diagonal right edge of the green block
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
                                  Inches(GW - 1.05), Inches(0),
                                  Inches(1.05), Inches(GH))
    tri.fill.solid()
    tri.fill.fore_color.rgb = WHITE
    tri.line.fill.background()
    # flipH=1 rotates the right-angle to the correct corner
    xfrm = tri._element.spPr.xfrm
    if xfrm is None:
        xfrm = etree.SubElement(tri._element.spPr, qn("a:xfrm"))
    xfrm.set("flipH", "1")

    # Logo image inside green block
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH,
                                  Inches(0.2), Inches(0.25),
                                  Inches(2.3), Inches(1.1))
    else:
        tb = add_tb(slide, 0.25, 0.3, GW - 0.5, 0.7, wrap=False)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = "SOCIOSCOPE"
        fnt(r, 20, bold=True, color=WHITE, face="Arial")
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = "F  O  O  D"
        fnt(r2, 12, bold=True, color=WHITE, face="Arial")

    # "FOOD" sub-label below logo (only shown when logo image is used)
    if os.path.exists(LOGO_PATH):
        tb2 = add_tb(slide, 0.25, 1.4, GW - 0.6, 0.5, wrap=False)
        p = tb2.text_frame.paragraphs[0]
        r = p.add_run(); r.text = "F  O  O  D"
        fnt(r, 13, bold=True, color=WHITE, face="Arial")

    # ── top-right thin green accent ──────────────────────────────────────────
    add_rect(slide, W - 2.0, 0, 2.0, 0.32, GREEN)

    # ── "CERTIFICATE" ────────────────────────────────────────────────────────
    tb = add_tb(slide, 0.5, 1.55, W - 1.0, 1.25, wrap=False)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "CERTIFICATE"
    fnt(r, 62, bold=True, color=DARK_GREEN, face="Calibri")
    set_char_spacing(r, 250)   # ~2.5 pt extra spacing

    # ── "Proudly presented to :" ─────────────────────────────────────────────
    tb = add_tb(slide, 0.5, 2.85, 7.0, 0.45, wrap=False)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Proudly presented to :"
    fnt(r, 13, color=BODY_CLR, face="Calibri")

    # ── horizontal line 1 ────────────────────────────────────────────────────
    add_hline(slide, 0.5, 3.42, W - 1.0)

    # ── initiative name ───────────────────────────────────────────────────────
    n = len(name)
    fs = 34 if n <= 20 else 28 if n <= 35 else 24 if n <= 50 else 20

    tb = add_tb(slide, 0.5, 3.47, W - 1.0, 1.1)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name.upper()
    fnt(r, fs, bold=True, color=DARK_GREEN, face="Calibri")

    # ── horizontal line 2 ────────────────────────────────────────────────────
    add_hline(slide, 0.5, 4.68, W - 1.0)

    # ── body text ─────────────────────────────────────────────────────────────
    tb = add_tb(slide, 0.9, 4.82, W - 1.8, 1.1)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    parts = [
        ("Chosen for your exceptional ", False),
        ("innovation", True),
        (" and ", False),
        ("impact", True),
        (" in advancing sustainable food systems, we express our gratitude "
         "for embodying The Food Socioscope's values and ", False),
        ("inspiring", True),
        (" global change.", False),
    ]
    for text, bold in parts:
        r = p.add_run(); r.text = text
        fnt(r, 14, bold=bold, color=BODY_CLR, face="Calibri")

    # ── signatures ────────────────────────────────────────────────────────────
    SIG_Y = 5.9

    # Left — Saadi Lahlou
    tb = add_tb(slide, 0.5, SIG_Y + 0.45, 4.0, 1.3, wrap=False)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "SAADI LAHLOU"
    fnt(r, 11, bold=True, color=BODY_CLR, face="Calibri")
    for line in ["Professor and Director", "Paris Institute for Advanced Study"]:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = line
        fnt(r2, 10, italic=True, color=BODY_CLR, face="Calibri")

    # Right — Helga Nowotny
    tb = add_tb(slide, W - 5.0, SIG_Y + 0.45, 4.5, 1.5, wrap=False)
    tf = tb.text_frame
    for i, line in enumerate([
        "HELGA NOWOTNY",
        "Professor and Chair of the Science Advisory Board",
        "Strategic Advisory Board",
        "Complexity Science Hub Vienna",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = line
        fnt(r, 11 if i == 0 else 10,
            bold=(i == 0), italic=(i > 0), color=BODY_CLR, face="Calibri")

    # ── bottom logos (text placeholders — replace with images as needed) ──────
    tb = add_tb(slide, 3.3, 7.25, 5.2, 0.7)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "IEA PARIS          Complexity Science Hub          NOMIS FOUNDATION"
    fnt(r, 8, color=BODY_CLR, face="Calibri")

    # ── page URL (small, bottom centre) ──────────────────────────────────────
    tb = add_tb(slide, 0.5, H - 0.35, W - 1.0, 0.3, wrap=False)
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = url
    fnt(r, 8, color=RGBColor(0x88, 0x88, 0x88), face="Calibri")

# ── build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)

for init in initiatives:
    make_slide(prs, init)

prs.save(OUTPUT_PATH)
print(f"OK  {len(initiatives)} certificates saved to:\n   {OUTPUT_PATH}")
